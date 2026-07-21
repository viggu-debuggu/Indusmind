import os
import io
import zipfile
import csv
from typing import List, Dict, Any, Tuple
import fitz  # PyMuPDF
import docx
import openpyxl
import pptx
from PIL import Image

from app.core.logging import logger
from app.ai.ocr import perform_ocr_on_pil_image, is_tesseract_available


def clean_extracted_text(text: str) -> str:
    """Removes excessive spaces, invalid characters, and normalizes lines."""
    if not text:
        return ""
    # Remove null characters
    text = text.replace("\x00", "")
    lines = [line.strip() for line in text.split("\n")]
    # Remove empty lines, join lines with space
    cleaned = " ".join([line for line in lines if line])
    # Replace multiple spaces with a single space
    import re
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


class DocumentProcessor:
    """Orchestrates structured text extraction per page/section from multiple file formats."""

    @classmethod
    def process_file(cls, file_bytes: bytes, filename: str, mime_type: str) -> List[Dict[str, Any]]:
        """Determines file type and routes to the appropriate extraction method."""
        ext = os.path.splitext(filename)[1].lower()
        logger.info("processing_document_file", filename=filename, extension=ext, mime=mime_type)

        try:
            if ext == ".pdf":
                return cls._extract_pdf(file_bytes)
            elif ext in (".docx", ".doc"):
                return cls._extract_docx(file_bytes)
            elif ext in (".xlsx", ".xls"):
                return cls._extract_xlsx(file_bytes)
            elif ext == ".pptx":
                return cls._extract_pptx(file_bytes)
            elif ext == ".csv":
                return cls._extract_csv(file_bytes)
            elif ext == ".txt":
                return cls._extract_txt(file_bytes)
            elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
                return cls._extract_image(file_bytes)
            elif ext == ".zip":
                return cls._extract_zip(file_bytes)
            else:
                # Default fallback read as txt
                return cls._extract_txt(file_bytes)
        except Exception as e:
            logger.exception("document_processing_failed", filename=filename, error=str(e))
            raise ValueError(f"Failed to process document {filename}: {str(e)}")

    @classmethod
    def _extract_pdf(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Extracts text page-by-page. Fallbacks to PyTesseract OCR if page is scanned (empty text)."""
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages_data = []

        tesseract_available = is_tesseract_available()

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            # Detect section heading (simple heuristics on font sizes or capitals)
            section = "Introduction"
            blocks = page.get_text("blocks")
            for b in blocks:
                # block: (x0, y0, x1, y1, "text", block_no, block_type)
                b_text = b[4].strip()
                if b_text and len(b_text) < 100 and b_text.isupper():
                    section = b_text
                    break

            cleaned = clean_extracted_text(text)

            # If page text is very short/empty and OCR is available, trigger page rendering OCR
            if len(cleaned) < 50 and tesseract_available:
                logger.info("scanned_pdf_page_detected_running_ocr", page=page_num + 1)
                try:
                    pix = page.get_pixmap(dpi=150)
                    image_data = pix.tobytes("png")
                    image = Image.open(io.BytesIO(image_data))
                    ocr_text = perform_ocr_on_pil_image(image)
                    cleaned = clean_extracted_text(ocr_text)
                    section = "Scanned Page OCR"
                except Exception as ocr_err:
                    logger.error("pdf_page_ocr_failed", page=page_num + 1, error=str(ocr_err))

            pages_data.append({
                "text": cleaned,
                "page": page_num + 1,
                "section": section,
                "metadata": {"type": "PDF"}
            })

        return pages_data

    @classmethod
    def _extract_docx(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Extracts structured headings, paragraphs, lists, and tables via python-docx."""
        doc_stream = io.BytesIO(file_bytes)
        doc = docx.Document(doc_stream)
        
        pages_data = []
        current_section = "General"
        current_text = []
        
        # Word documents don't have hard physical pages easily mapped,
        # we will segment text by heading blocks (logical pages/sections).
        logical_page = 1

        for elem in doc.element.body:
            if elem.tag.endswith('p'):
                p = docx.text.paragraph.Paragraph(elem, doc)
                p_text = p.text.strip()
                if not p_text:
                    continue
                
                # Check style names for headings
                style_name = p.style.name.lower()
                if "heading" in style_name:
                    # Save former logical page context
                    if current_text:
                        pages_data.append({
                            "text": clean_extracted_text(" ".join(current_text)),
                            "page": logical_page,
                            "section": current_section,
                            "metadata": {"type": "DOCX"}
                        })
                        logical_page += 1
                        current_text = []
                    current_section = p_text
                else:
                    current_text.append(p_text)
                    
            elif elem.tag.endswith('tbl'):
                table = docx.table.Table(elem, doc)
                # Formulate table as text representation
                table_lines = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_lines.append(" | ".join(row_data))
                if table_lines:
                    current_text.append("\nTable Data:\n" + "\n".join(table_lines) + "\n")

        # Save remaining chunk
        if current_text or current_section:
            pages_data.append({
                "text": clean_extracted_text(" ".join(current_text)),
                "page": logical_page,
                "section": current_section,
                "metadata": {"type": "DOCX"}
            })

        return pages_data

    @classmethod
    def _extract_xlsx(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Extracts spreadsheet sheets and cell values into structured text grids."""
        xlsx_stream = io.BytesIO(file_bytes)
        wb = openpyxl.load_workbook(xlsx_stream, read_only=True, data_only=True)
        
        pages_data = []
        
        # Map sheets as logical sections/pages
        for page_idx, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            sheet_rows = []
            
            # Read first 1000 rows to prevent memory blowout
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count > 1000:
                    break
                # Filter out completely empty rows
                row_str = [str(val).strip() for val in row if val is not None]
                if row_str:
                    sheet_rows.append(", ".join(row_str))
                row_count += 1
                
            text = f"Sheet: {sheet_name}\n" + "\n".join(sheet_rows)
            pages_data.append({
                "text": clean_extracted_text(text),
                "page": page_idx + 1,
                "section": f"Sheet: {sheet_name}",
                "metadata": {"type": "XLSX"}
            })
            
        return pages_data

    @classmethod
    def _extract_pptx(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Extracts slides and titles structure via python-pptx."""
        pptx_stream = io.BytesIO(file_bytes)
        prs = pptx.Presentation(pptx_stream)
        pages_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_title = f"Slide {slide_idx + 1}"
            slide_texts = []
            
            # Check if slide has header placeholder
            if slide.shapes.title:
                slide_title = slide.shapes.title.text.strip()
                
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if text != slide_title:
                        slide_texts.append(text)
                        
            cleaned = clean_extracted_text(" ".join(slide_texts))
            pages_data.append({
                "text": cleaned,
                "page": slide_idx + 1,
                "section": slide_title,
                "metadata": {"type": "PPTX"}
            })
            
        return pages_data

    @classmethod
    def _extract_csv(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Reads CSV sheet rows."""
        csv_text = file_bytes.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(csv_text))
        rows_str = []
        for row in reader:
            if row:
                rows_str.append(" | ".join(row))
                
        # Split CSV values into logical pages of 100 rows each
        pages_data = []
        chunk_size = 100
        for i in range(0, len(rows_str), chunk_size):
            chunk = rows_str[i:i + chunk_size]
            text = "\n".join(chunk)
            page_num = (i // chunk_size) + 1
            pages_data.append({
                "text": clean_extracted_text(text),
                "page": page_num,
                "section": f"Data Block {page_num}",
                "metadata": {"type": "CSV"}
            })
            
        return pages_data

    @classmethod
    def _extract_txt(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Reads plain text file contents, chunked into logical 3000-char blocks."""
        raw_text = file_bytes.decode("utf-8", errors="ignore")
        cleaned_raw = raw_text.replace("\r\n", "\n")
        
        pages_data = []
        char_limit = 3000
        
        # Split by logical paragraph boundaries or fixed length characters
        paragraphs = cleaned_raw.split("\n\n")
        current_page_text = []
        current_char_count = 0
        page_num = 1
        
        for p in paragraphs:
            p = p.strip()
            if not p:
                continue
            current_page_text.append(p)
            current_char_count += len(p)
            
            if current_char_count >= char_limit:
                pages_data.append({
                    "text": clean_extracted_text(" ".join(current_page_text)),
                    "page": page_num,
                    "section": f"Section {page_num}",
                    "metadata": {"type": "TXT"}
                })
                page_num += 1
                current_page_text = []
                current_char_count = 0
                
        if current_page_text:
            pages_data.append({
                "text": clean_extracted_text(" ".join(current_page_text)),
                "page": page_num,
                "section": f"Section {page_num}",
                "metadata": {"type": "TXT"}
            })
            
        return pages_data

    @classmethod
    def _extract_image(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Uses Tesseract to compute characters from raw image formats (PNG, JPG, etc.)."""
        tesseract_available = is_tesseract_available()
        
        if not tesseract_available:
            logger.warning("tesseract_ocr_unavailable_aborting_image_extract")
            return [{
                "text": "Tesseract OCR binary not found. Image content could not be extracted.",
                "page": 1,
                "section": "Image Scanned",
                "metadata": {"type": "Image", "status": "Failed OCR"}
            }]
            
        try:
            image = Image.open(io.BytesIO(file_bytes))
            ocr_text = perform_ocr_on_pil_image(image)
            return [{
                "text": clean_extracted_text(ocr_text),
                "page": 1,
                "section": "Image OCR",
                "metadata": {"type": "Image"}
            }]
        except Exception as e:
            logger.error("image_extraction_failed", error=str(e))
            return []

    @classmethod
    def _extract_zip(cls, file_bytes: bytes) -> List[Dict[str, Any]]:
        """Unpacks zip archives and processes all files inside recursively."""
        zip_stream = io.BytesIO(file_bytes)
        pages_data = []
        
        with zipfile.ZipFile(zip_stream) as archive:
            for file_info in archive.infolist():
                # Skip directories and hidden system files
                if file_info.is_dir() or os.path.basename(file_info.filename).startswith('.'):
                    continue
                    
                filename = file_info.filename
                mime_type = "application/octet-stream"  # generic fallback
                
                try:
                    nested_bytes = archive.read(file_info.filename)
                    # Process nested file
                    nested_pages = cls.process_file(nested_bytes, filename, mime_type)
                    # Annotate nested filenames in page metadata
                    for page in nested_pages:
                        page["section"] = f"Archive: {filename} > {page['section']}"
                        page["metadata"]["archive_filename"] = filename
                        pages_data.append(page)
                except Exception as nested_err:
                    logger.error("archive_file_extraction_item_failed", filename=filename, error=str(nested_err))
                    continue
                    
        return pages_data
